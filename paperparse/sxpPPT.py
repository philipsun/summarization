#-------------------------------------------------------------------------------
# Name:        sxpPPT
# Purpose:
#
# Author:      sunxp
#
# Created:     01-04-2015
# Copyright:   (c) sunxp 2015
# Licence:     <your licence>
#-------------------------------------------------------------------------------
#code=utf-8
from pptx import Presentation
from pptx.util import Inches, Pt
import time
import sxpParseXHMLPaper

def MakePPTFromText(sxptext):
    prs = Presentation()

    pptname = sxptext.fname + '.pptx'

    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    t = time.time()
    title.text = sxptext.title
    subtitle.text = "Xiaoping Sun's Python programm made!\n"+time.strftime('%Y-%d-%m',t)
#**************************
    bullet_slide_layout = prs.slide_layouts[1]

    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Topic List:'

    section_title_set = sxptext.section_title_set
    tf = body_shape.text_frame

    fi = 1
    for sec in section_title_set:
        if sec.t_type == 'ltx_title_document':
            continue
        if sec.level<=1:
            level = 0;
        else:
            level = sec.level
        p = tf.add_paragraph()
        p.text = sec.title#'Use _TextFrame.text for first bullet'
        p.level = level
        p.font.size = Pt(24)
        fi = fi + 1;
        if fi % 10 == 0:
            slide = prs.slides.add_slide(bullet_slide_layout)
            shapes = slide.shapes

            title_shape = shapes.title
            body_shape = shapes.placeholders[1]

            title_shape.text = 'Topic List:'

            section_title_set = sxptext.section_title_set
            tf = body_shape.text_frame

##        p = tf.add_paragraph()
##        p.text = 'Use _TextFrame.add_paragraph() for subsequent bullets'
##        p.level = 2

#***************************
    prs.save('test.pptx')

def MakePPTFromTextA(sxptext):
    prs = Presentation()

    pptname = sxptext.fname + '.pptx'

    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = sxptext.title
    t = time.localtime()
    subtitle.text = "Xiaoping Sun's Python program made!\n"+time.strftime('%Y-%d-%m',t)
#************************** Add section titles to Overview
    bullet_slide_layout = prs.slide_layouts[1]

    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Overview:'

    section_title_set = sxptext.section_title_set
    tf = body_shape.text_frame

    fi = 1
    for sec in section_title_set:
        if sec.t_type == 'ltx_title_document':
            continue
        if sec.t_type == 'ltx_title_abstract':
            continue
        if sec.level<=1:
            level = 0;
            if fi == 1:
                tf.text = sec.title
                fi = fi + 1;
            else:
                p = tf.add_paragraph()
                p.text = sec.title#'Use _TextFrame.text for first bullet'
                p.level = level
                #p.font.size = Pt(24)
                fi = fi + 1;
                if fi % 10 == 0:
                    slide = prs.slides.add_slide(bullet_slide_layout)
                    shapes = slide.shapes

                    title_shape = shapes.title
                    body_shape = shapes.placeholders[1]

                    title_shape.text = 'Topic List:'

                    section_title_set = sxptext.section_title_set
                    tf = body_shape.text_frame
        else:
            level = sec.level
            continue
##        p = tf.add_paragraph()
##        p.text = 'Use _TextFrame.add_paragraph() for subsequent bullets'
##        p.level = 2

#***************************Add section content for each paragraph

#***************************
    prs.save('test.pptx')

def TestMakeHelloWorld():
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Hello, World!"
    subtitle.text = "python-pptx was here!"

    prs.save('test.pptx')
def TestPPT():

    prs = Presentation()
    bullet_slide_layout = prs.slide_layouts[1]

    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Adding a Bullet Slide'

    tf = body_shape.text_frame
    tf.text = 'Find the bullet slide layout\n hello and this is a test'

    p = tf.add_paragraph()
    p.text = 'Use _TextFrame.text for first bullet'
    p.level = 0

    p = tf.add_paragraph()
    p.text = 'Use _TextFrame.add_paragraph() for subsequent bullets'
    p.level = 1

    prs.save('test.pptx')
def main():
    TestPPT()

if __name__ == '__main__':
    main()
